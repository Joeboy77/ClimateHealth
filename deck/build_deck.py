"""Build the ClimaHealth Predict pitch deck.

Regenerate:  cd backend && uv run python ../deck/build_deck.py

Twelve slides. White ground, one accent, one alarm, a serif reserved for figures
and verdicts. Built in code so every number on the deck comes from one place.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

DECK = Path(__file__).resolve().parent
SHOTS = DECK / "screenshots"
OUTPUT = DECK / "ClimaHealth-Predict.pptx"

WIDTH, HEIGHT = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1A, 0x1D, 0x1B)
MUTED = RGBColor(0x6F, 0x6D, 0x66)
FAINT = RGBColor(0xA8, 0xA5, 0x9E)
ACCENT = RGBColor(0x0E, 0x6E, 0x63)
ALARM = RGBColor(0xA3, 0x21, 0x18)
RULE = RGBColor(0xE4, 0xE2, 0xDC)
TINT = RGBColor(0xEF, 0xF5, 0xF3)
CANVAS = RGBColor(0xFA, 0xF8, 0xF5)

DISPLAY = "Georgia"
BODY = "Helvetica Neue"

MARGIN = Inches(0.85)
CONTENT_WIDTH = WIDTH - MARGIN * 2


def text_box(slide, left, top, width, height, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.paragraphs[0].alignment = align
    return frame


def write(
    frame,
    text,
    size,
    colour=INK,
    font=BODY,
    bold=False,
    space_after=0,
    line=None,
    spacing=None,
    first=False,
    align=None,
):
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.space_after = Pt(space_after)
    if line is not None:
        paragraph.line_spacing = line
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = colour
    run.font.name = font
    run.font.bold = bold
    if spacing is not None:
        run.font._rPr.set("spc", str(int(spacing * 100)))
    return paragraph


def rect(slide, left, top, width, height, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def hairline(slide, top, left=MARGIN, width=CONTENT_WIDTH, colour=RULE):
    return rect(slide, left, top, width, Emu(9525), fill=colour)


def blank(presentation, number=None, dark=False):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK if dark else RGBColor(0xFF, 0xFF, 0xFF)

    # A short accent stroke in the top margin: enough to make the deck feel set
    # rather than typed, without becoming decoration that competes with the data.
    rect(slide, MARGIN, Inches(0.34), Inches(0.62), Emu(28575), fill=ACCENT)

    if number is not None:
        frame = text_box(slide, WIDTH - Inches(1.5), HEIGHT - Inches(0.62), Inches(0.65), Inches(0.3))
        write(frame, f"{number:02d}", 10, FAINT, BODY, True, spacing=1.2, first=True, align=PP_ALIGN.RIGHT)
        mark = text_box(slide, MARGIN, HEIGHT - Inches(0.62), Inches(4.0), Inches(0.3))
        write(mark, "CLIMAHEALTH PREDICT", 9, FAINT, BODY, True, spacing=1.8, first=True)
    return slide


def eyebrow(slide, text, top=Inches(0.66)):
    frame = text_box(slide, MARGIN, top, CONTENT_WIDTH, Inches(0.3))
    write(frame, text.upper(), 11, ACCENT, BODY, True, spacing=1.8, first=True)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def panel(slide, left, top, width, height, fill=TINT):
    return rect(slide, left, top, width, height, fill=fill)


def shot(slide, name, left, top, max_height, slot, caption=None, framed=True):
    """Place a screenshot inside a slot, fitting both its width and its height so it
    can never spill onto a neighbour or off the canvas. A missing file leaves a
    labelled frame, so the gap is visible rather than silently absent."""
    path = SHOTS / name
    inner = Emu(int(slot * 0.92))

    if path.exists():
        added = slide.shapes.add_picture(str(path), left, top, width=inner)
        if added.height > max_height:
            scale = max_height / added.height
            added.width = Emu(int(added.width * scale))
            added.height = max_height
        added.left = left + Emu(int((slot - added.width) / 2))
        if framed:
            rect(slide, added.left, added.top, added.width, added.height, line=RULE)
        bottom = added.top + added.height
    else:
        width = Emu(int(inner))
        height = Emu(int(min(max_height, width * 2)))
        shp = rect(
            slide, left + Emu(int((slot - width) / 2)), top, width, height,
            fill=CANVAS, line=RULE,
        )
        label = shp.text_frame
        label.word_wrap = True
        label.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(label, f"[ {name} ]", 10, FAINT, BODY, True, first=True, align=PP_ALIGN.CENTER)
        bottom = top + height

    if caption is not None:
        cap = text_box(slide, left, bottom + Inches(0.12), slot, Inches(0.45))
        write(cap, caption, 11, MUTED, BODY, False, first=True, align=PP_ALIGN.CENTER)
    return bottom


def figure_rail(slide, items, top, left=MARGIN, gap=Inches(3.05), size=30, colour=ACCENT):
    for index, (value, label) in enumerate(items):
        column = text_box(slide, left + gap * index, top, gap - Inches(0.25), Inches(1.1))
        write(column, value, size, colour, DISPLAY, False, 2, first=True)
        write(column, label, 12, MUTED, BODY, False, 0, line=1.3)


# ---------------------------------------------------------------- slides


def title_slide(presentation):
    slide = blank(presentation)

    panel(slide, Inches(0), Inches(0), Inches(0.28), HEIGHT, fill=ACCENT)

    frame = text_box(slide, Inches(1.3), Inches(2.05), Inches(11.0), Inches(2.8))
    write(frame, "CLIMAHEALTH PREDICT", 12, ACCENT, BODY, True, 20, spacing=2.4, first=True)
    write(frame, "Ghana knows the weather.", 44, INK, DISPLAY, False, 2, line=1.06)
    write(frame, "It does not yet know what the weather", 44, INK, DISPLAY, False, 2, line=1.06)
    write(frame, "is about to do to people.", 44, ACCENT, DISPLAY, False, 0, line=1.06)

    hairline(slide, Inches(5.25), left=Inches(1.3), width=Inches(11.0))
    footer = text_box(slide, Inches(1.3), Inches(5.5), Inches(11.0), Inches(1.0))
    write(
        footer,
        "A climate-health early-warning platform. It turns today's climate into a ranked, "
        "explained health warning for the weeks ahead.",
        15,
        MUTED,
        BODY,
        False,
        8,
        first=True,
    )
    write(footer, "GREENRES HACKATHON 2026", 11, FAINT, BODY, True, spacing=1.8)
    notes(
        slide,
        "Open with the one sentence, then stop. Say: every agency in this country "
        "already has weather data. Nobody turns it into who gets sick, where, and when. "
        "That gap is the whole product.",
    )
    return slide


def problem_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "The problem")

    frame = text_box(slide, MARGIN, Inches(1.05), Inches(11.6), Inches(0.9))
    write(frame, "Climate-driven disease is Ghana's largest health burden.", 30, INK, DISPLAY, False, 8, first=True)
    write(
        frame,
        "Ghana is one of the 15 highest-burden malaria countries in the world. "
        "Cholera, meningitis and heat illness track the weather just as closely.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.4,
    )

    panel(slide, MARGIN, Inches(2.35), Inches(5.6), Inches(2.0), fill=CANVAS)
    left_frame = text_box(slide, Inches(1.15), Inches(2.6), Inches(5.0), Inches(1.6))
    write(left_frame, "6.7 million", 36, ALARM, DISPLAY, False, 4, first=True)
    write(left_frame, "malaria cases in Ghana in 2024, and 11,635 deaths.", 15, INK, BODY, True, 4, line=1.3)
    write(left_frame, "WHO, World Malaria Report 2024", 11, FAINT, BODY, False, 0)

    right_frame = text_box(slide, Inches(7.1), Inches(2.3), Inches(5.4), Inches(0.4))
    write(right_frame, "AND THE CASES ARRIVE LATE", 11, MUTED, BODY, True, spacing=1.6, first=True)

    lags = [
        ("Malaria", "2 to 6 weeks", "after rainfall and standing water"),
        ("Cholera", "1 to 3 weeks", "after flooding and unsafe water"),
        ("Meningitis", "2 to 8 weeks", "after dry, dusty harmattan air"),
    ]
    top = Inches(2.8)
    for condition, window, cause in lags:
        row = text_box(slide, Inches(7.1), top, Inches(5.4), Inches(0.9))
        write(row, f"{condition}   {window}", 17, ACCENT, DISPLAY, False, 2, first=True)
        write(row, cause, 12, MUTED, BODY, False, 0)
        top += Inches(0.85)

    hairline(slide, Inches(5.55))
    closing = text_box(slide, MARGIN, Inches(5.8), CONTENT_WIDTH, Inches(1.2))
    write(
        closing,
        "That delay is not the problem. It is the opportunity.",
        20,
        INK,
        DISPLAY,
        False,
        8,
        first=True,
    )
    write(
        closing,
        "It is the only window in which a district can still act before people fall ill. "
        "Ghana currently reacts when the cases appear, and the climate data that precedes "
        "them is free, public and updated daily.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.45,
    )
    notes(
        slide,
        "Lead with the number, not the adjective. 6.7 million cases, 11,635 deaths, "
        "one year, one country. Then turn it: all of it follows weather we can already "
        "see, and it arrives weeks late. That lag is the product.",
    )
    return slide


def why_act_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "Why anybody acts")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(11.6), Inches(0.9))
    write(frame, "Nobody changes their day for a weather forecast.", 32, INK, DISPLAY, False, 8, first=True)
    write(
        frame,
        "Ghana already has weather warnings. People read them and carry on, because rain "
        "is not news and a forecast asks nothing of you.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.4,
    )

    panel(slide, MARGIN, Inches(2.35), Inches(5.5), Inches(3.05), fill=CANVAS)
    left_frame = text_box(slide, Inches(1.15), Inches(2.6), Inches(4.9), Inches(2.6))
    write(left_frame, "A WEATHER SERVICE SAYS", 10, MUTED, BODY, True, 10, spacing=1.6, first=True)
    write(left_frame, "“Heavy rain expected Thursday.”", 19, MUTED, DISPLAY, False, 12, line=1.2)
    write(left_frame, "Severity:  it rains every year", 13, MUTED, BODY, False, 3)
    write(left_frame, "About me:  not obviously", 13, MUTED, BODY, False, 3)
    write(left_frame, "To do:  nothing is asked", 13, MUTED, BODY, False, 10)
    write(left_frame, "Nothing changes.", 14, MUTED, BODY, True, 0)

    panel(slide, Inches(6.75), Inches(2.35), Inches(5.75), Inches(3.05), fill=TINT)
    right_frame = text_box(slide, Inches(7.05), Inches(2.6), Inches(5.15), Inches(2.6))
    write(right_frame, "CLIMAHEALTH SAYS", 10, ACCENT, BODY, True, 10, spacing=1.6, first=True)
    write(
        right_frame,
        "“Malaria risk is rising here. Cases in 2 to 6 weeks. "
        "Children under five and pregnant women most at risk. "
        "Empty standing water tonight.”",
        16,
        INK,
        DISPLAY,
        False,
        12,
        line=1.25,
    )
    write(right_frame, "Severity:  a named disease", 13, INK, BODY, False, 3)
    write(right_frame, "About me:  my district, my children", 13, INK, BODY, False, 3)
    write(right_frame, "To do:  one thing, tonight, free", 13, INK, BODY, False, 10)
    write(right_frame, "Something to actually do.", 14, ACCENT, BODY, True, 0)

    close = text_box(slide, MARGIN, Inches(5.6), CONTENT_WIDTH, Inches(1.25))
    write(
        close,
        "We do not warn people about the weather. We warn them about what it is about "
        "to do to their children, and give them one thing to do about it.",
        18,
        INK,
        DISPLAY,
        False,
        8,
        first=True,
        line=1.25,
    )
    write(
        close,
        "Protection Motivation Theory: protective action needs severity, personal "
        "vulnerability, and a response the person believes they can carry out. Fear on "
        "its own produces avoidance, which is why every warning here ends in an action.",
        12,
        FAINT,
        BODY,
        False,
        0,
        line=1.35,
    )
    notes(
        slide,
        "This is your strongest slide. Read the left column flatly, the right column "
        "with weight. Then the line: we do not warn people about the weather, we warn "
        "them about what it is about to do to their children. "
        "If a judge calls it fear-mongering, the answer is already on the slide: fear "
        "without an action backfires, so every warning ends in one thing to do tonight.",
    )
    return slide


def solution_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "The solution")

    frame = text_box(slide, MARGIN, Inches(1.05), Inches(11.6), Inches(1.3))
    write(frame, "ClimaHealth Predict", 32, INK, DISPLAY, False, 6, first=True)
    write(
        frame,
        "Reads a district's climate today and says which health risks are rising, why, "
        "how long before cases appear, and who is most affected.",
        16,
        MUTED,
        BODY,
        False,
        0,
        line=1.4,
    )

    parts = [
        ("For the public", "Dawuro", "Forecast and weather, hazard reporting with a photo, and a daily lesson that pays in NHIS cover."),
        ("For the agencies", "Command Platform", "260 districts ranked and explained, alerts, an incident room, readiness, a view per mandate."),
        ("For everyone else", "USSD and SMS", "The same warning on any handset, with no data and no app, because a smartphone is not a requirement."),
    ]
    left = MARGIN
    for audience, name, detail in parts:
        panel(slide, left, Inches(2.75), Inches(3.7), Inches(2.35), fill=CANVAS)
        column = text_box(slide, left + Inches(0.3), Inches(3.0), Inches(3.1), Inches(2.0))
        write(column, audience.upper(), 10, MUTED, BODY, True, 8, spacing=1.6, first=True)
        write(column, name, 21, ACCENT, DISPLAY, False, 8)
        write(column, detail, 13, MUTED, BODY, False, 0, line=1.45)
        left += Inches(3.95)

    hairline(slide, Inches(5.5))
    frame = text_box(slide, MARGIN, Inches(5.75), CONTENT_WIDTH, Inches(1.2))
    write(
        frame,
        "The brain is a rules engine, not a guess.",
        20,
        INK,
        DISPLAY,
        False,
        8,
        first=True,
    )
    write(
        frame,
        "Published epidemiological thresholds, evaluated the same way every time. Same "
        "inputs, same output, and every warning can name the conditions that caused it. "
        "AI never decides risk; it only phrases what the engine already decided, in "
        "English or Twi. That is what makes this defensible to a health ministry.",
        14,
        MUTED,
        BODY,
        False,
        0,
        line=1.45,
    )
    notes(
        slide,
        "Three doors, one engine. Then the bottom line, which is your credibility: "
        "deterministic rules from published thresholds, reproducible and explainable. "
        "The AI is a translator, not a decider. Say that sentence exactly.",
    )
    return slide


def dashboard_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "The agency platform")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(11.6), Inches(0.9))
    write(frame, "260 districts, ranked and explained, every day.", 30, INK, DISPLAY, False, 8, first=True)
    write(
        frame,
        "Live engine output, not a mock. Each agency opens on the layer it is "
        "responsible for: health risk for GHS, dust and PM10 for EPA, rainfall for NADMO.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.4,
    )

    shot(slide, "w1-national.png", MARGIN, Inches(2.45), Inches(3.5), Inches(5.6),
         caption="Ghana Health Service · 51 of 260 districts at high risk or above")
    shot(slide, "w3-epa.png", Inches(6.9), Inches(2.45), Inches(3.5), Inches(5.6),
         caption="Environmental Protection Agency · dust and PM10 layer active")

    notes(
        slide,
        "Point at the 51 of 260, then at the map, then at the second screenshot and say: "
        "same engine, same districts, different question. An air quality officer should "
        "not have to hunt for air quality.",
    )
    return slide


def loop_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "Institutional integration")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(11.6), Inches(0.9))
    write(frame, "A report is a claim until somebody goes and looks.", 30, INK, DISPLAY, False, 8, first=True)
    write(
        frame,
        "Ɔhwɛfoɔ is an on-ground officer role in the platform. Validation and repair "
        "are different jobs held by different people, and neither can take the other's step.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.4,
    )

    stages = [
        ("Citizen reports", "photo, location, note", "Dawuro app"),
        ("Ɔhwɛfoɔ validates", "goes to the site, confirms", "25% → 50%"),
        ("Agency works", "GHS, NADMO, EPA, Assembly", "50% → 75%"),
        ("Resolved", "closed with what was done", "100%"),
    ]
    left = MARGIN
    for index, (name, detail, meta) in enumerate(stages):
        panel(slide, left, Inches(2.6), Inches(2.75), Inches(1.95), fill=TINT if index == 1 else CANVAS)
        column = text_box(slide, left + Inches(0.25), Inches(2.8), Inches(2.3), Inches(1.6))
        write(column, f"0{index + 1}", 12, ACCENT, DISPLAY, True, 6, first=True)
        write(column, name, 16, INK, BODY, True, 4)
        write(column, detail, 12, MUTED, BODY, False, 6, line=1.3)
        write(column, meta, 12, ACCENT, BODY, True)
        left += Inches(3.0)

    hairline(slide, Inches(4.95))
    frame = text_box(slide, MARGIN, Inches(5.2), CONTENT_WIDTH, Inches(1.6))
    write(
        frame,
        "Every step is written down: who moved it, when, and what they found.",
        19,
        INK,
        DISPLAY,
        False,
        8,
        first=True,
    )
    write(
        frame,
        "The person who filed the report watches the same bar the officers do. Reporting "
        "a hazard into a form that never answers is how people learn not to bother. "
        "The agencies do not have to trust the app; they have to trust their own officer.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.45,
    )
    notes(
        slide,
        "Judges said institutional integration was unproven. This is the answer: a named "
        "role, a permission model, an append-only trail, and the citizen sees it close. "
        "End on the line about trusting their own officer.",
    )
    return slide


def app_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "The app, and who it is for")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(5.7), Inches(1.9))
    write(frame, "Dawuro", 30, INK, DISPLAY, False, 6, first=True)
    write(
        frame,
        "Named after the gong-gong that carried public warnings in Ghana for centuries. "
        "The people this must reach are the least likely to own a good phone, read "
        "fluently, or speak English at home. Each of those is designed for.",
        14,
        MUTED,
        BODY,
        False,
        0,
        line=1.45,
    )

    rows = [
        ("Five languages", "English, Twi, Ga, Ewe, Dagbani"),
        ("Read aloud", "Forecasts, lessons and answers can be spoken"),
        ("Never a wrong voice", "No voice for a language, it says so and stays silent"),
        ("Age-banded", "A child and a grandmother get different wording"),
        ("Works offline", "Last forecast kept; reports queue and send themselves"),
    ]
    top = Inches(3.25)
    for name, detail in rows:
        row = text_box(slide, MARGIN, top, Inches(5.7), Inches(0.66))
        write(row, name, 14, ACCENT, BODY, True, 2, first=True)
        write(row, detail, 12, MUTED, BODY, False, 0, line=1.3)
        top += Inches(0.68)

    shots = [
        ("m1-home.png", "Today's risk"),
        ("m2-quiz.png", "Learn by answering"),
        ("m5-twi.png", "The same warning in Twi"),
    ]
    left = Inches(6.75)
    for name, caption in shots:
        shot(slide, name, left, Inches(1.15), Inches(5.0), slot=Inches(1.9), caption=caption)
        left += Inches(1.95)

    notes(
        slide,
        "Twenty seconds. Point at the Twi screen last: it tells the reader that their "
        "phone cannot read Twi aloud yet and that the wording is awaiting review by a "
        "Twi speaker. Nobody else in this room will show a screen admitting its limits.",
    )
    return slide


def ussd_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "Reach")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(5.8), Inches(1.4))
    write(frame, "A smartphone is not a requirement.", 30, INK, DISPLAY, False, 8, first=True)
    write(
        frame,
        "Ghana has more mobile connections than people, yet three in ten Ghanaians do "
        "not use the internet at all. An app alone never reaches them.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.45,
    )

    panel(slide, MARGIN, Inches(2.95), Inches(5.8), Inches(1.5), fill=CANVAS)
    figure_rail(
        slide,
        [("38.3m", "mobile connections"),
         ("69.9%", "use the internet")],
        Inches(3.2),
        left=Inches(1.15),
        gap=Inches(2.75),
        size=26,
    )
    source = text_box(slide, Inches(1.15), Inches(4.62), Inches(5.2), Inches(0.3))
    write(source, "GSMA Intelligence / DataReportal, Digital 2025 Ghana", 10, FAINT, BODY, False, first=True)

    frame = text_box(slide, MARGIN, Inches(5.15), Inches(5.8), Inches(1.8))
    write(frame, "Live on Africa's Talking.", 19, INK, DISPLAY, False, 8, first=True)
    write(
        frame,
        "Dial the shortcode, choose a language, choose a region, choose a district, get "
        "today's warning. No data, no app, no smartphone. SMS pushes the same warning to "
        "a district when risk crosses the threshold.",
        14,
        MUTED,
        BODY,
        False,
        0,
        line=1.45,
    )

    shots = [
        ("u1-region.png", "Choose region"),
        ("u2-district.png", "Choose district"),
        ("u3-warning.png", "Today's warning"),
    ]
    left = Inches(7.0)
    for name, caption in shots:
        shot(slide, name, left, Inches(1.15), Inches(5.05), slot=Inches(1.8), caption=caption)
        left += Inches(1.85)

    notes(
        slide,
        "Judges flagged USSD and SMS as not implemented. These are real screenshots of "
        "a real dial-through on Africa's Talking. Say the three-in-ten number, then "
        "point at the last screen: that is the same engine output, on a phone with no "
        "internet.",
    )
    return slide


def reward_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "Rewards that are fundable")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(11.6), Inches(0.9))
    write(frame, "Points buy health cover, not cash.", 30, INK, DISPLAY, False, 8, first=True)
    write(
        frame,
        "3,500 points is one year of NHIS cover. Derived, not chosen: the adult premium "
        "runs GHS 7.20 to 48 by income band, about GHS 30 to 50 for informal workers, "
        "and we price a point at one pesewa.",
        15,
        MUTED,
        BODY,
        False,
        0,
        line=1.4,
    )

    shot(slide, "w2-renewals.png", MARGIN, Inches(2.5), Inches(3.25), Inches(6.1),
         caption="The renewal queue · GHS only · seeded demonstration Guardians")

    panel(slide, Inches(7.35), Inches(2.5), Inches(5.15), Inches(3.3), fill=TINT)
    right = text_box(slide, Inches(7.7), Inches(2.75), Inches(4.55), Inches(2.9))
    write(right, "No money leaves the platform.", 17, INK, BODY, True, 8, first=True, line=1.25)
    write(
        right,
        "No float to hold, no payout licence, no transfer that cannot be recalled. "
        "The platform never says a renewal happened: it records who earned one and "
        "hands that to Ghana Health Service, who renew and confirm.",
        13,
        MUTED,
        BODY,
        False,
        12,
        line=1.45,
    )
    write(right, "About 44% of Ghanaians hold no active cover.", 15, ACCENT, BODY, True, 6, line=1.3)
    write(
        right,
        "They are the same people most exposed to climate-driven illness. The reward "
        "removes the exact cost that kept them out, which is why they come back "
        "tomorrow. The streak forgives a missed day, and wrong answers still earn, so "
        "nobody is locked out of health information for getting one wrong.",
        13,
        MUTED,
        BODY,
        False,
        0,
        line=1.45,
    )

    notes(
        slide,
        "Two answers on one slide. First: cash removed, nothing to fund or regulate. "
        "Second, which is the adoption question every judge asks: 44% have no active "
        "cover, and the loop pays in exactly the thing they are going without. "
        "That is the wedge, not a gimmick.",
    )
    return slide


def answers_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "What you asked us last time")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(11.6), Inches(0.7))
    write(frame, "Three criticisms, three answers.", 30, INK, DISPLAY, False, 0, first=True)

    rows = [
        ("Field validation unproven",
         "Ɔhwɛfoɔ: an officer stands where the report was filed and confirms it before "
         "any agency spends a truck. Every validation signed and timestamped."),
        ("Institutional integration unproven",
         "Reports move through named agency stages with a permission model, and every "
         "agency gets a view scoped to its own mandate."),
        ("Reward feasibility unproven",
         "Cash removed. Points buy NHIS cover priced against the real premium band, "
         "issued by Ghana Health Service, not by us."),
    ]
    top = Inches(2.0)
    for criticism, answer in rows:
        hairline(slide, top - Inches(0.2))
        row = text_box(slide, MARGIN, top, Inches(4.0), Inches(1.0))
        write(row, criticism, 16, ALARM, BODY, True, 0, line=1.3, first=True)
        response = text_box(slide, Inches(5.2), top, Inches(7.3), Inches(1.0))
        write(response, answer, 14, INK, BODY, False, 0, line=1.45, first=True)
        top += Inches(1.32)

    panel(slide, MARGIN, Inches(5.95), CONTENT_WIDTH, Inches(1.05), fill=TINT)
    ask = text_box(slide, Inches(1.15), Inches(6.1), Inches(11.0), Inches(0.85))
    write(ask, "WHAT WE ARE ASKING FOR", 10, ACCENT, BODY, True, 6, spacing=1.6, first=True)
    write(
        ask,
        "One district and one season. Give us Madina and a Ghana Health Service link, and "
        "we will run the engine against DHIMS2 case records to show how early the warning "
        "actually was. It is the only part we cannot build ourselves.",
        14,
        INK,
        BODY,
        False,
        0,
        line=1.4,
    )
    notes(
        slide,
        "Say each criticism out loud before its answer. Judges remember being listened "
        "to. Then the ask, and stop talking. Naming the one thing you cannot do alone "
        "reads as confidence and gives them something concrete to say yes to.",
    )
    return slide


def close_slide(presentation, number):
    slide = blank(presentation, number, dark=True)

    frame = text_box(slide, Inches(1.3), Inches(2.3), Inches(11.0), Inches(2.4))
    write(frame, "CLIMAHEALTH PREDICT", 12, RGBColor(0x5F, 0xB5, 0xA8), BODY, True, 20, spacing=2.4, first=True)
    write(frame, "The climate signal is already there.", 36, RGBColor(0xFF, 0xFF, 0xFF), DISPLAY, False, 2, line=1.08)
    write(frame, "We turn it into a warning", 36, RGBColor(0xFF, 0xFF, 0xFF), DISPLAY, False, 2, line=1.08)
    write(frame, "somebody can act on.", 36, RGBColor(0x5F, 0xB5, 0xA8), DISPLAY, False, 0, line=1.08)

    rect(slide, Inches(1.3), Inches(5.2), Inches(11.0), Emu(9525), fill=RGBColor(0x3A, 0x3F, 0x3C))
    items = [
        ("260", "districts, daily"),
        ("5", "disease pathways"),
        ("3", "ways in"),
        ("963", "tests passing"),
    ]
    for index, (value, label) in enumerate(items):
        column = text_box(slide, Inches(1.3) + Inches(2.75) * index, Inches(5.45), Inches(2.5), Inches(1.0))
        write(column, value, 28, RGBColor(0x5F, 0xB5, 0xA8), DISPLAY, False, 2, first=True)
        write(column, label, 12, RGBColor(0x9A, 0x9E, 0x9A), BODY, False, 0)

    notes(
        slide,
        "Close on the three lines and stop. Let the numbers sit. Backup slide has every "
        "source if they ask where a figure came from.",
    )
    return slide


def sources_slide(presentation, number):
    slide = blank(presentation, number)
    eyebrow(slide, "Backup · sources and limits")

    frame = text_box(slide, MARGIN, Inches(1.0), Inches(11.6), Inches(0.6))
    write(frame, "Where the numbers come from.", 28, INK, DISPLAY, False, 0, first=True)

    rows = [
        ("Malaria burden", "6.7m cases, 11,635 deaths in Ghana, 2024", "WHO, World Malaria Report 2024"),
        ("NHIS coverage", "18.5m active members, about 56% of the population", "National Health Insurance Authority, 2025"),
        ("NHIS premium", "GHS 7.20 to 48 a year; GHS 30 to 50 informal", "NHIA scheme documentation"),
        ("Mobile reach", "38.3m connections; 69.9% use the internet", "GSMA Intelligence / DataReportal 2025"),
        ("Behaviour change", "Severity, vulnerability and a doable action", "Protection Motivation Theory, Rogers"),
        ("Climate data", "Daily observations for all 260 districts", "Open-Meteo, live in the product"),
        ("Boundaries", "260 district polygons on the map", "geoBoundaries, CC BY 4.0"),
    ]
    top = Inches(1.85)
    for topic, claim, source in rows:
        hairline(slide, top - Inches(0.1))
        write(text_box(slide, MARGIN, top, Inches(2.8), Inches(0.55)), topic, 13, INK, BODY, True, 0, first=True)
        write(text_box(slide, Inches(3.8), top, Inches(5.2), Inches(0.55)), claim, 12, MUTED, BODY, False, 0, line=1.3, first=True)
        write(text_box(slide, Inches(9.3), top, Inches(3.3), Inches(0.55)), source, 11, ACCENT, BODY, False, 0, line=1.3, first=True)
        top += Inches(0.66)

    panel(slide, MARGIN, Inches(6.5), CONTENT_WIDTH, Inches(0.72), fill=CANVAS)
    limit = text_box(slide, Inches(1.15), Inches(6.62), Inches(11.0), Inches(0.55))
    write(
        limit,
        "Known limit, stated plainly: the engine has not yet been backtested against "
        "DHIMS2 case data. That is the ask, and the one thing between a defensible "
        "model and a validated one.",
        13,
        MUTED,
        BODY,
        False,
        first=True,
        line=1.35,
    )
    notes(
        slide,
        "Do not present this. It exists so that when a judge asks where a number came "
        "from you turn one page and answer. Volunteering the DHIMS2 gap before being "
        "asked reads as confident; being caught without it reads as careless.",
    )
    return slide


def build() -> Path:
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = WIDTH, HEIGHT

    title_slide(presentation)
    problem_slide(presentation, 2)
    why_act_slide(presentation, 3)
    solution_slide(presentation, 4)
    dashboard_slide(presentation, 5)
    loop_slide(presentation, 6)
    app_slide(presentation, 7)
    ussd_slide(presentation, 8)
    reward_slide(presentation, 9)
    answers_slide(presentation, 10)
    close_slide(presentation, 11)
    sources_slide(presentation, 12)

    presentation.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"{path}  ({path.stat().st_size / 1024:.0f} KB)")
